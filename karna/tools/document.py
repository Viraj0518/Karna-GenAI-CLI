"""Document tool -- reads and extracts content from PDF, Word, PowerPoint,
Excel, and CSV files.

Returns content as formatted markdown text. Supports page ranges for PDFs
and truncates very long documents (max 50K chars).

Optional dependencies: pdfplumber, python-docx, python-pptx, openpyxl.
Install with: pip install karna[docs]
"""

from __future__ import annotations

import csv
import os
from pathlib import Path
from typing import Any

from karna.security import is_safe_path
from karna.tools.base import BaseTool

_MAX_CHARS = 50_000

# --------------------------------------------------------------------------- #
#  Optional dependency imports
# --------------------------------------------------------------------------- #

try:
    import pdfplumber  # type: ignore[import-untyped]

    _HAS_PDFPLUMBER = True
except ModuleNotFoundError:
    pdfplumber = None  # type: ignore[assignment]
    _HAS_PDFPLUMBER = False

try:
    import docx as _docx  # type: ignore[import-untyped]

    _HAS_DOCX = True
except ModuleNotFoundError:
    _docx = None  # type: ignore[assignment]
    _HAS_DOCX = False

try:
    import pptx as _pptx  # type: ignore[import-untyped]

    _HAS_PPTX = True
except ModuleNotFoundError:
    _pptx = None  # type: ignore[assignment]
    _HAS_PPTX = False

try:
    import openpyxl  # type: ignore[import-untyped]

    _HAS_OPENPYXL = True
except ModuleNotFoundError:
    openpyxl = None  # type: ignore[assignment]
    _HAS_OPENPYXL = False

_INSTALL_HINT = "Install with: pip install karna[docs]"

# Supported extensions
_SUPPORTED_EXTENSIONS = frozenset({".pdf", ".docx", ".pptx", ".xlsx", ".csv"})


# --------------------------------------------------------------------------- #
#  Format-specific extractors
# --------------------------------------------------------------------------- #


def _parse_page_range(pages: str, total_pages: int) -> list[int]:
    """Parse a page range string into a list of 0-based page indices.

    Supports formats like '3', '1-5', '10-20'.
    """
    pages = pages.strip()
    if "-" in pages:
        parts = pages.split("-", 1)
        start = int(parts[0]) - 1
        end = int(parts[1])  # inclusive, 1-based
        start = max(0, start)
        end = min(end, total_pages)
        return list(range(start, end))
    else:
        idx = int(pages) - 1
        if 0 <= idx < total_pages:
            return [idx]
        return []


def _extract_pdf(path: Path, pages: str | None) -> str:
    """Extract text and tables from a PDF file using pdfplumber."""
    if not _HAS_PDFPLUMBER:
        return f"[error] pdfplumber is not installed. {_INSTALL_HINT}"

    parts: list[str] = []
    with pdfplumber.open(path) as pdf:
        total = len(pdf.pages)
        parts.append(f"# PDF: {path.name} ({total} pages)\n")

        if pages:
            indices = _parse_page_range(pages, total)
            if not indices:
                return f"[error] Invalid page range '{pages}' for PDF with {total} pages."
        else:
            indices = list(range(total))

        for idx in indices:
            page = pdf.pages[idx]
            parts.append(f"## Page {idx + 1}\n")

            text = page.extract_text()
            if text:
                parts.append(text)

            tables = page.extract_tables()
            for ti, table in enumerate(tables):
                parts.append(f"\n**Table {ti + 1}:**\n")
                parts.append(_format_table(table))

            parts.append("")

    return "\n".join(parts)


def _extract_docx(path: Path) -> str:
    """Extract paragraphs and tables from a Word document."""
    if not _HAS_DOCX:
        return f"[error] python-docx is not installed. {_INSTALL_HINT}"

    doc = _docx.Document(str(path))
    parts: list[str] = [f"# Word Document: {path.name}\n"]

    for para in doc.paragraphs:
        text = para.text.strip()
        if text:
            # Preserve heading styles
            if para.style and para.style.name and para.style.name.startswith("Heading"):
                try:
                    level = int(para.style.name.split()[-1])
                except (ValueError, IndexError):
                    level = 2
                parts.append(f"{'#' * (level + 1)} {text}\n")
            else:
                parts.append(text)

    for ti, table in enumerate(doc.tables):
        parts.append(f"\n**Table {ti + 1}:**\n")
        rows = []
        for row in table.rows:
            rows.append([cell.text.strip() for cell in row.cells])
        parts.append(_format_table(rows))

    return "\n".join(parts)


def _extract_pptx(path: Path) -> str:
    """Extract text from PowerPoint slides."""
    if not _HAS_PPTX:
        return f"[error] python-pptx is not installed. {_INSTALL_HINT}"

    prs = _pptx.Presentation(str(path))
    parts: list[str] = [f"# PowerPoint: {path.name} ({len(prs.slides)} slides)\n"]

    for si, slide in enumerate(prs.slides, 1):
        parts.append(f"## Slide {si}\n")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        parts.append(text)
        parts.append("")

    return "\n".join(parts)


def _extract_xlsx(path: Path) -> str:
    """Extract sheet data from an Excel workbook."""
    if not _HAS_OPENPYXL:
        return f"[error] openpyxl is not installed. {_INSTALL_HINT}"

    wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
    parts: list[str] = [f"# Excel: {path.name} ({len(wb.sheetnames)} sheets)\n"]

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        parts.append(f"## Sheet: {sheet_name}\n")
        rows: list[list[str]] = []
        for row in ws.iter_rows(values_only=True):
            rows.append([str(cell) if cell is not None else "" for cell in row])
        if rows:
            parts.append(_format_table(rows))
        else:
            parts.append("(empty sheet)")
        parts.append("")

    wb.close()
    return "\n".join(parts)


def _extract_csv(path: Path) -> str:
    """Extract data from a CSV file."""
    parts: list[str] = [f"# CSV: {path.name}\n"]

    with open(path, newline="", encoding="utf-8", errors="replace") as fh:
        reader = csv.reader(fh)
        rows: list[list[str]] = []
        for row in reader:
            rows.append(row)

    if rows:
        parts.append(_format_table(rows))
    else:
        parts.append("(empty file)")

    return "\n".join(parts)


# --------------------------------------------------------------------------- #
#  Helpers
# --------------------------------------------------------------------------- #


def _format_table(rows: list[list[Any]]) -> str:
    """Format a list of rows as a markdown table."""
    if not rows:
        return "(empty table)"

    # Normalise: ensure all rows have same column count
    max_cols = max(len(r) for r in rows)
    normalised = []
    for row in rows:
        cells = [(str(c).strip() if c is not None else "") for c in row]
        while len(cells) < max_cols:
            cells.append("")
        normalised.append(cells)

    # Build markdown table
    lines: list[str] = []
    header = normalised[0]
    lines.append("| " + " | ".join(header) + " |")
    lines.append("| " + " | ".join("---" for _ in header) + " |")
    for row in normalised[1:]:
        lines.append("| " + " | ".join(row) + " |")

    return "\n".join(lines)


def _truncate(text: str) -> str:
    """Truncate text to _MAX_CHARS if it exceeds the limit."""
    if len(text) <= _MAX_CHARS:
        return text
    return text[:_MAX_CHARS] + f"\n\n... [truncated — {len(text)} total chars, showing first {_MAX_CHARS}]"


# --------------------------------------------------------------------------- #
#  DocumentTool
# --------------------------------------------------------------------------- #


class DocumentTool(BaseTool):
    """Read and extract content from PDF, Word, PowerPoint, Excel, and CSV files."""

    name = "document"
    description = "Read and extract content from PDF, Word, PowerPoint, and Excel files"
    sequential = False
    parameters: dict[str, Any] = {
        "type": "object",
        "properties": {
            "file_path": {
                "type": "string",
                "description": "Path to the document",
            },
            "pages": {
                "type": "string",
                "description": "Page range for PDFs (e.g., '1-5', '3', '10-20')",
            },
        },
        "required": ["file_path"],
    }

    async def execute(self, **kwargs: Any) -> str:
        """Extract content from the specified document file."""
        file_path_str: str = kwargs.get("file_path", "")
        pages: str | None = kwargs.get("pages")

        if not file_path_str:
            return "[error] No file_path provided."

        path = Path(os.path.expanduser(file_path_str)).resolve()

        if not is_safe_path(file_path_str):
            return "[error] Blocked: path points to a sensitive or disallowed location."

        if not path.exists():
            return f"[error] File not found: {path}"
        if not path.is_file():
            return f"[error] Not a file: {path}"

        ext = path.suffix.lower()
        if ext not in _SUPPORTED_EXTENSIONS:
            supported = ", ".join(sorted(_SUPPORTED_EXTENSIONS))
            return f"[error] Unsupported document format: '{ext}'. Supported: {supported}"

        try:
            if ext == ".pdf":
                content = _extract_pdf(path, pages)
            elif ext == ".docx":
                content = _extract_docx(path)
            elif ext == ".pptx":
                content = _extract_pptx(path)
            elif ext == ".xlsx":
                content = _extract_xlsx(path)
            elif ext == ".csv":
                content = _extract_csv(path)
            else:
                return f"[error] Unsupported document format: '{ext}'"
        except Exception as exc:  # noqa: BLE001
            return f"[error] Failed to process {path.name}: {exc}"

        return _truncate(content)
